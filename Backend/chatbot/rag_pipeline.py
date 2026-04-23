"""
RAG Pipeline for Chatbot
Handles retrieval-augmented generation for timetable, schedules, and attendance topics
"""
import os
import glob
import base64
from datetime import datetime
from typing import List, Dict, Optional, Tuple
import hashlib
from pathlib import Path
try:
    from langchain_community.embeddings import HuggingFaceEmbeddings
    from langchain_community.vectorstores import Chroma
    try:
        from langchain_text_splitters import RecursiveCharacterTextSplitter
    except ImportError:
        from langchain.text_splitter import RecursiveCharacterTextSplitter
    # RetrievalQA is deprecated in LangChain 1.x, we'll use create_retrieval_chain instead
    RetrievalQA = None  # Will use alternative approach
    try:
        from langchain.chains.combine_documents import create_stuff_documents_chain
        from langchain.chains import create_retrieval_chain
        HAS_RETRIEVAL_CHAIN = True
    except ImportError:
        HAS_RETRIEVAL_CHAIN = False
        try:
            # Fallback to old RetrievalQA if available
            from langchain.chains.retrieval_qa.base import RetrievalQA
        except ImportError:
            try:
                from langchain.chains import RetrievalQA
            except ImportError:
                RetrievalQA = None
    try:
        from langchain_core.prompts import PromptTemplate
    except ImportError:
        from langchain.prompts import PromptTemplate
    try:
        from langchain_core.documents import Document
    except ImportError:
        from langchain.schema import Document
    import chromadb
    from chromadb.config import Settings
    
    # Try different Groq import paths
    try:
        from langchain_groq import ChatGroq
        Groq = ChatGroq  # Use ChatGroq as Groq
    except ImportError:
        try:
            from langchain_community.llms import Groq
        except ImportError:
            try:
                from langchain.llms import Groq
            except ImportError:
                Groq = None
                print("WARNING: Chatbot: Groq LLM not available in LangChain")
except ImportError:
    # Fallback for older langchain versions
    try:
        from langchain.embeddings import HuggingFaceEmbeddings
        from langchain.vectorstores import Chroma
        from langchain.text_splitter import RecursiveCharacterTextSplitter
        from langchain.chains import RetrievalQA
        from langchain.prompts import PromptTemplate
        from langchain.schema import Document
        import chromadb
        from chromadb.config import Settings
        
        try:
            from langchain.llms import Groq
        except ImportError:
            Groq = None
    except ImportError:
        print("WARNING: Chatbot: LangChain or ChromaDB not installed")
        HuggingFaceEmbeddings = None
        Chroma = None
        RecursiveCharacterTextSplitter = None
        RetrievalQA = None
        Groq = None
        PromptTemplate = None
        Document = None
        chromadb = None
        Settings = None
import json
import traceback
import re

try:
    import requests
except ImportError:
    requests = None
try:
    import pytesseract
except ImportError:
    pytesseract = None
try:
    from PIL import Image, ImageOps, ImageFilter
except ImportError:
    Image = None
try:
    from pypdf import PdfReader
except Exception:
    try:
        from PyPDF2 import PdfReader
    except Exception:
        PdfReader = None
try:
    from docx import Document as DocxDocument
except Exception:
    DocxDocument = None
try:
    from pptx import Presentation as PptxPresentation
except Exception:
    PptxPresentation = None

# Initialize embedding model (reuse the same model as schedules)
EMBEDDING_MODEL_NAME = "BAAI/bge-base-en-v1.5"
# Groq model id (override with GROQ_CHAT_MODEL in .env).
# Use a stronger default model for better timetable Q&A quality.
GROQ_CHAT_MODEL = os.getenv("GROQ_CHAT_MODEL", "llama-3.3-70b-versatile")
GROQ_CHAT_COMPLETIONS_URL = "https://api.groq.com/openai/v1/chat/completions"
GROQ_TIMETABLE_VISION_MODEL = os.getenv(
    "GROQ_TIMETABLE_VISION_MODEL",
    "meta-llama/llama-4-scout-17b-16e-instruct",
)
GROQ_TIMETABLE_VISION_FALLBACK_MODELS = [
    "meta-llama/llama-4-scout-17b-16e-instruct",
    "meta-llama/llama-4-maverick-17b-128e-instruct",
    # Keep non-prefixed variants as last-resort compatibility fallbacks.
    "llama-4-scout-17b-16e-instruct",
    "llama-4-maverick-17b-128e-instruct",
]
# Student vector matching behavior (Chroma returns distance; lower is better for cosine distance)
STUDENT_RAG_MAX_DISTANCE = float(os.getenv("STUDENT_RAG_MAX_DISTANCE", "0.45"))
STUDENT_RELATED_RAG_MAX_DISTANCE = float(os.getenv("STUDENT_RELATED_RAG_MAX_DISTANCE", "0.65"))
ADMIN_STUDENT_QUERY_MAX_DOCS = int(os.getenv("ADMIN_STUDENT_QUERY_MAX_DOCS", "25"))
BACKEND_PUBLIC_BASE_URL = (os.getenv("BACKEND_PUBLIC_BASE_URL", "http://127.0.0.1:5000") or "http://127.0.0.1:5000").rstrip("/")
# /course mode: Chroma distance below this => use retrieved chunks; else inject full raw file text (truncated).
COURSE_MATERIAL_RAG_MAX_DISTANCE = float(os.getenv("COURSE_MATERIAL_RAG_MAX_DISTANCE", "0.45"))
COURSE_MATERIAL_RELATED_MAX_DISTANCE = float(os.getenv("COURSE_MATERIAL_RELATED_MAX_DISTANCE", "0.80"))
COURSE_MATERIAL_FULL_CONTEXT_MAX_CHARS = int(os.getenv("COURSE_MATERIAL_FULL_CONTEXT_MAX_CHARS", "120000"))
MAX_COURSE_MATERIAL_FILES = 6
VALID_COURSE_KEYS = frozenset({"computerScience", "chemistry", "physics"})
COURSE_MATERIALS_COLLECTION = "course_materials"

class RAGPipeline:
    def __init__(self, chroma_db_path: str = None):
        """Initialize RAG pipeline with embedding model and ChromaDB"""
        self.embeddings = None
        self.timetable_store = None
        self.schedules_store = None
        self.attendance_store = None
        self.llm = None
        self._groq_api_key = None
        self.chroma_client = None
        self._timetable_context_cache = ""
        self._timetable_context_cache_source = ""
        self._course_material_store_cache = {}
        _backend_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        self._course_materials_upload_root = os.path.join(_backend_dir, "uploads", "course_materials")
        self.chroma_db_path = chroma_db_path or os.path.join(
            _backend_dir,
            'chroma_db'
        )
        # Create ChromaDB folder immediately
        try:
            os.makedirs(self.chroma_db_path, exist_ok=True)
            print(f"✅ Chatbot: ChromaDB folder created/verified at {self.chroma_db_path}")
        except Exception as e:
            print(f"⚠️  Chatbot: Failed to create ChromaDB folder: {e}")
        self._ensure_subject_chroma_dirs()
        
        self._initialize_chromadb()
        self._initialize_embeddings()
        self._initialize_llm()

    def _build_attendance_collection_name(self, user_email: str, user_role: str, user_course: Optional[str]) -> str:
        """Create user-scoped attendance collection to avoid cross-user data leakage."""
        raw = f"{user_role}:{user_course or 'none'}:{(user_email or '').lower()}"
        digest = hashlib.sha1(raw.encode("utf-8")).hexdigest()[:16]
        return f"attendance_{digest}"
    
    def _initialize_chromadb(self):
        """Initialize ChromaDB client for local storage"""
        try:
            if chromadb is None:
                print("WARNING: Chatbot: ChromaDB not available - install with: pip install chromadb")
                return
            
            # Ensure folder exists
            os.makedirs(self.chroma_db_path, exist_ok=True)
            
            # Create ChromaDB client with local persistence
            if Settings:
                self.chroma_client = chromadb.PersistentClient(
                    path=self.chroma_db_path,
                    settings=Settings(
                        anonymized_telemetry=False,
                        allow_reset=True
                    )
                )
            else:
                # Fallback for older ChromaDB versions
                self.chroma_client = chromadb.PersistentClient(path=self.chroma_db_path)
            
            print(f"✅ Chatbot: ChromaDB initialized at {self.chroma_db_path}")
            print(f"✅ Chatbot: ChromaDB folder exists: {os.path.exists(self.chroma_db_path)}")
            
            # Try to load existing collections
            self._load_existing_collections()
        except Exception as e:
            print(f"⚠️  Chatbot: Failed to initialize ChromaDB: {e}")
            import traceback
            traceback.print_exc()
            self.chroma_client = None
    
    def _load_existing_collections(self):
        """Load existing ChromaDB collections if they exist"""
        try:
            if not self.chroma_client or not self.embeddings or Chroma is None:
                return
            
            # Try to load timetable collection
            try:
                collection = self.chroma_client.get_collection(name="timetable")
                if collection.count() > 0:
                    self.timetable_store = Chroma(
                        client=self.chroma_client,
                        collection_name="timetable",
                        embedding_function=self.embeddings
                    )
                    print(f"✅ Chatbot: Loaded existing timetable collection ({collection.count()} documents)")
            except:
                pass
            
            # Try to load schedules collection
            try:
                collection = self.chroma_client.get_collection(name="schedules")
                if collection.count() > 0:
                    self.schedules_store = Chroma(
                        client=self.chroma_client,
                        collection_name="schedules",
                        embedding_function=self.embeddings
                    )
                    print(f"✅ Chatbot: Loaded existing schedules collection ({collection.count()} documents)")
            except:
                pass
            
            # Attendance collection is loaded dynamically per user, so we don't load it here
        except Exception as e:
            print(f"⚠️  Chatbot: Error loading existing collections: {e}")
    
    def _initialize_embeddings(self):
        """Initialize HuggingFace embeddings"""
        try:
            if HuggingFaceEmbeddings is None:
                print("⚠️  Chatbot: HuggingFaceEmbeddings not available")
                return
            self.embeddings = HuggingFaceEmbeddings(
                model_name=EMBEDDING_MODEL_NAME,
                model_kwargs={'device': 'cpu'},
                encode_kwargs={'normalize_embeddings': True}
            )
            print("✅ Chatbot: Embedding model initialized")
        except Exception as e:
            print(f"⚠️  Chatbot: Failed to initialize embeddings: {e}")
            self.embeddings = None
    
    def _initialize_llm(self):
        """Initialize Groq LLM (LangChain) and cache API key for HTTP fallback."""
        self._groq_api_key = (os.getenv("GROQ_API_KEY") or "").strip() or None
        if not self._groq_api_key:
            print("WARNING: Chatbot: GROQ_API_KEY not found in environment")
            return
        try:
            if Groq is None:
                print("WARNING: Chatbot: Groq LLM package not available — will use Groq HTTP API if needed")
                return
            # Check if it's ChatGroq (newer) or Groq (older)
            try:
                from langchain_groq import ChatGroq
                # api_key alias matches langchain-groq ChatGroq; use same model as HTTP path
                self.llm = ChatGroq(
                    model=GROQ_CHAT_MODEL,
                    api_key=self._groq_api_key,
                    temperature=0.2,
                    max_tokens=900,
                )
            except ImportError:
                # Fallback to old Groq API
                self.llm = Groq(
                    model_name=GROQ_CHAT_MODEL,
                    groq_api_key=self._groq_api_key,
                    temperature=0.2,
                    max_tokens=900,
                )
            print("✅ Chatbot: LLM initialized")
        except Exception as e:
            print(f"⚠️  Chatbot: Failed to initialize LangChain Groq (HTTP fallback still available): {e}")
            import traceback
            traceback.print_exc()
            self.llm = None

    def _groq_http_complete(self, prompt: str, max_tokens: int = 512, timeout: int = 90) -> Optional[str]:
        """Direct Groq OpenAI-compatible API (same approach as working quiz code in auth/routes.py)."""
        if not self._groq_api_key or requests is None:
            return None
        try:
            response = requests.post(
                GROQ_CHAT_COMPLETIONS_URL,
                headers={
                    "Authorization": f"Bearer {self._groq_api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": GROQ_CHAT_MODEL,
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.2,
                    "max_tokens": max_tokens,
                },
                timeout=timeout,
            )
            if response.status_code != 200:
                print(f"⚠️  Chatbot: Groq HTTP {response.status_code}: {response.text[:800]}")
                return None
            data = response.json()
            text = data["choices"][0]["message"]["content"]
            return (text or "").strip()
        except Exception as e:
            print(f"⚠️  Chatbot: Groq HTTP request failed: {e}")
            traceback.print_exc()
            return None

    def _groq_http_complete_with_image(self, prompt: str, image_path: str, max_tokens: int = 900, timeout: int = 120) -> Optional[str]:
        """Send timetable image directly to Groq Vision model."""
        if not self._groq_api_key or requests is None:
            return None
        try:
            suffix = Path(image_path).suffix.lower()
            if suffix not in {".png", ".jpg", ".jpeg"}:
                return None
            mime = "image/png" if suffix == ".png" else "image/jpeg"

            with open(image_path, "rb") as f:
                data_uri = f"data:{mime};base64,{base64.b64encode(f.read()).decode('utf-8')}"

            def normalize_model_id(model_id: str) -> str:
                m = (model_id or "").strip()
                # Fix common typos seen in runtime/env values
                m = m.replace("insttruct", "instruct").replace("--", "-")
                return m

            models_to_try = [normalize_model_id(GROQ_TIMETABLE_VISION_MODEL)] + [
                m for m in GROQ_TIMETABLE_VISION_FALLBACK_MODELS if m != GROQ_TIMETABLE_VISION_MODEL
            ]
            models_to_try = [normalize_model_id(m) for m in models_to_try if normalize_model_id(m)]

            # de-duplicate preserving order
            deduped = []
            seen = set()
            for m in models_to_try:
                if m in seen:
                    continue
                seen.add(m)
                deduped.append(m)
            models_to_try = deduped

            for model_name in models_to_try:
                response = requests.post(
                    GROQ_CHAT_COMPLETIONS_URL,
                    headers={
                        "Authorization": f"Bearer {self._groq_api_key}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": model_name,
                        "messages": [
                            {
                                "role": "user",
                                "content": [
                                    {"type": "text", "text": prompt},
                                    {"type": "image_url", "image_url": {"url": data_uri}},
                                ],
                            }
                        ],
                        "temperature": 0.2,
                        "max_tokens": max_tokens,
                    },
                    timeout=timeout,
                )

                if response.status_code == 200:
                    data = response.json()
                    return (data["choices"][0]["message"]["content"] or "").strip()

                body = (response.text or "")[:1200]
                print(f"⚠️  Chatbot: Groq Vision ({model_name}) HTTP {response.status_code}: {body[:800]}")
                lower_body = body.lower()
                # Retry automatically on decommissioned/invalid model cases.
                if "decommissioned" in lower_body or "model_decommissioned" in lower_body or "does not exist" in lower_body:
                    continue
                # For non-model errors, stop retry loop.
                return None

            return None
        except Exception as e:
            print(f"⚠️  Chatbot: Groq Vision request failed: {e}")
            traceback.print_exc()
            return None

    def _get_latest_timetable_upload_path(self) -> str:
        """Get latest clear uploaded timetable file path (ignore debug/processed artifacts)."""
        try:
            def is_preferred_timetable_file(path: str) -> bool:
                name = os.path.basename(path).lower()
                # Ignore OCR/debug artifacts so LLM always receives the original clear file.
                if "_processed" in name or "processed" in name:
                    return False
                if "_debug" in name or "debug" in name:
                    return False
                return True

            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            uploads_candidates = [
                os.path.join(base_dir, "uploads", "timetables"),
                os.path.join(base_dir, "timetable", "uploads"),
            ]
            files = []
            for folder in uploads_candidates:
                if not folder or not os.path.exists(folder):
                    continue
                for ext in ("*.png", "*.jpg", "*.jpeg", "*.pdf"):
                    files.extend(glob.glob(os.path.join(folder, ext)))
            if not files:
                return ""

            preferred = [f for f in files if is_preferred_timetable_file(f)]
            candidates = preferred if preferred else files
            # Prefer latest; on ties prefer larger file (usually clearer image quality).
            return max(candidates, key=lambda f: (os.path.getmtime(f), os.path.getsize(f)))
        except Exception:
            return ""

    def _invoke_llm_text(self, prompt: str) -> Optional[str]:
        """Invoke LangChain chat model with a proper message list (raw string invoke is unreliable for ChatGroq)."""
        if not self.llm:
            return None
        try:
            from langchain_core.messages import HumanMessage
            response = self.llm.invoke([HumanMessage(content=prompt)])
        except Exception as e1:
            try:
                response = self.llm.invoke(prompt)
            except Exception as e2:
                print(f"⚠️  Chatbot: LangChain invoke failed: {e1} | {e2}")
                return None
        if isinstance(response, str):
            return response.strip()
        content = getattr(response, "content", None)
        if content:
            return str(content).strip()
        return str(response).strip() if response is not None else None

    def _complete_prompt(self, prompt: str, max_tokens: int = 512) -> Optional[str]:
        """Prefer LangChain Groq; fall back to HTTP so answers always work when the key is valid."""
        text = self._invoke_llm_text(prompt)
        if text:
            return text
        return self._groq_http_complete(prompt, max_tokens=max_tokens)
    
    def load_timetable_embeddings(self, embeddings_folder: str):
        """Load timetable embeddings from files and store in ChromaDB"""
        try:
            if not self.embeddings or Document is None or Chroma is None or not self.chroma_client:
                print("⚠️  Chatbot: Required components not available")
                return False
            
            documents = []
            if os.path.exists(embeddings_folder):
                for filename in os.listdir(embeddings_folder):
                    if filename.endswith('_embeddings.txt'):
                        filepath = os.path.join(embeddings_folder, filename)
                        try:
                            with open(filepath, 'r', encoding='utf-8') as f:
                                content = f.read()
                                # Parse embeddings file - format: === key ===\nText: ...\nEmbedding: ...
                                current_text = None
                                for line in content.split('\n'):
                                    line = line.strip()
                                    if line.startswith('Text:'):
                                        current_text = line.replace('Text:', '').strip()
                                    elif line.startswith('===') and current_text:
                                        # Save previous text if exists
                                        if current_text and len(current_text) > 10:  # Minimum length
                                            documents.append(Document(
                                                page_content=current_text,
                                                metadata={"source": "timetable", "file": filename}
                                            ))
                                        current_text = None
                                    elif current_text and not line.startswith('Embedding:'):
                                        # Continue building text if it's a continuation
                                        if line:
                                            current_text += " " + line
                                
                                # Save last text if exists
                                if current_text and len(current_text) > 10:
                                    documents.append(Document(
                                        page_content=current_text,
                                        metadata={"source": "timetable", "file": filename}
                                    ))
                        except Exception as e:
                            print(f"⚠️  Error reading {filename}: {e}")
                            continue
            
            if documents and RecursiveCharacterTextSplitter:
                # Split documents into chunks
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=500,
                    chunk_overlap=50
                )
                texts = text_splitter.split_documents(documents)
                
                # Create or get ChromaDB collection for timetable
                collection_name = "timetable"
                try:
                    # Always rebuild from latest embedding files to keep vectors in sync.
                    self.chroma_client.delete_collection(name=collection_name)
                except Exception:
                    pass
                
                # Create new collection
                collection = self.chroma_client.create_collection(
                    name=collection_name,
                    metadata={"description": "Timetable embeddings for chatbot"}
                )
                
                # Store in ChromaDB via LangChain
                self.timetable_store = Chroma.from_documents(
                    documents=texts,
                    embedding=self.embeddings,
                    client=self.chroma_client,
                    collection_name=collection_name,
                    persist_directory=self.chroma_db_path
                )
                print(f"✅ Chatbot: Rebuilt timetable collection from {len(documents)} documents")
                return True
            else:
                if not documents:
                    print("⚠️  Chatbot: No timetable documents found")
                return False
        except Exception as e:
            print(f"⚠️  Chatbot: Error loading timetable embeddings: {e}")
            traceback.print_exc()
            return False
    
    def load_schedules_embeddings(self, embeddings_folder: str):
        """Load schedules embeddings from files and store in ChromaDB"""
        try:
            if not self.embeddings or Document is None or Chroma is None or not self.chroma_client:
                return False
            
            documents = []
            if os.path.exists(embeddings_folder):
                for filename in os.listdir(embeddings_folder):
                    if filename.endswith('_embeddings.txt'):
                        filepath = os.path.join(embeddings_folder, filename)
                        try:
                            with open(filepath, 'r', encoding='utf-8') as f:
                                content = f.read()
                                # Parse embeddings file - format: === key ===\nText: ...\nEmbedding: ...
                                current_text = None
                                for line in content.split('\n'):
                                    line = line.strip()
                                    if line.startswith('Text:'):
                                        current_text = line.replace('Text:', '').strip()
                                    elif line.startswith('===') and current_text:
                                        if current_text and len(current_text) > 10:
                                            documents.append(Document(
                                                page_content=current_text,
                                                metadata={"source": "schedule", "file": filename}
                                            ))
                                        current_text = None
                                    elif current_text and not line.startswith('Embedding:'):
                                        if line:
                                            current_text += " " + line
                                
                                if current_text and len(current_text) > 10:
                                    documents.append(Document(
                                        page_content=current_text,
                                        metadata={"source": "schedule", "file": filename}
                                    ))
                        except Exception as e:
                            print(f"⚠️  Error reading {filename}: {e}")
                            continue
            
            if documents and RecursiveCharacterTextSplitter:
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=500,
                    chunk_overlap=50
                )
                texts = text_splitter.split_documents(documents)
                
                # Create or get ChromaDB collection for schedules
                collection_name = "schedules"
                try:
                    # Always rebuild from latest embedding files to keep vectors in sync.
                    self.chroma_client.delete_collection(name=collection_name)
                except Exception:
                    pass
                
                collection = self.chroma_client.create_collection(
                    name=collection_name,
                    metadata={"description": "Schedule embeddings for chatbot"}
                )
                
                self.schedules_store = Chroma.from_documents(
                    documents=texts,
                    embedding=self.embeddings,
                    client=self.chroma_client,
                    collection_name=collection_name,
                    persist_directory=self.chroma_db_path
                )
                print(f"✅ Chatbot: Rebuilt schedules collection from {len(documents)} documents")
                return True
            else:
                if not documents:
                    print("⚠️  Chatbot: No schedule documents found")
                return False
        except Exception as e:
            print(f"⚠️  Chatbot: Error loading schedules embeddings: {e}")
            traceback.print_exc()
            return False

    def remove_vectors_for_file(self, source: str, upload_filename: Optional[str] = None, embeddings_filename: Optional[str] = None) -> bool:
        """
        Remove vectors from Chroma for a specific uploaded file.
        We store file metadata as '<upload_filename_with_dots_replaced>_embeddings.txt'.
        """
        try:
            if not self.chroma_client:
                return False

            if not embeddings_filename and upload_filename:
                embeddings_filename = upload_filename.replace(".", "_") + "_embeddings.txt"
            if not embeddings_filename:
                return False

            collection_name = "timetable" if source == "timetable" else "schedules" if source == "schedule" else None
            if not collection_name:
                return False

            collection = self.chroma_client.get_collection(name=collection_name)
            # Chroma metadata stores this as metadata.file
            collection.delete(where={"file": embeddings_filename})
            print(f"✅ Chatbot: Removed vectors from '{collection_name}' for file metadata: {embeddings_filename}")
            return True
        except Exception as e:
            print(f"⚠️  Chatbot: Failed removing vectors for {source}/{upload_filename}: {e}")
            return False

    def get_vector_counts(self) -> Dict:
        """Return current persisted vector counts for key collections."""
        counts = {"timetable": 0, "schedules": 0, "attendance": 0, "course_materials_by_course": {}}
        if not self.chroma_client:
            for ck in VALID_COURSE_KEYS:
                counts["course_materials_by_course"][ck] = self._course_material_chunk_count(ck)
            return counts
        try:
            counts["timetable"] = self.chroma_client.get_collection("timetable").count()
        except Exception:
            pass
        try:
            counts["schedules"] = self.chroma_client.get_collection("schedules").count()
        except Exception:
            pass
        # attendance collections are user-scoped and dynamic; this count is best-effort
        try:
            attendance_names = [
                c.name for c in self.chroma_client.list_collections()
                if getattr(c, "name", "").startswith("attendance_")
            ]
            total = 0
            for name in attendance_names:
                try:
                    total += self.chroma_client.get_collection(name).count()
                except Exception:
                    continue
            counts["attendance"] = total
        except Exception:
            pass
        for ck in VALID_COURSE_KEYS:
            counts["course_materials_by_course"][ck] = self._course_material_chunk_count(ck)
        return counts

    def _ensure_subject_chroma_dirs(self):
        """Physical folders under chroma_db/subjects/<courseKey> for per-subject vector stores."""
        try:
            root = os.path.join(self.chroma_db_path, "subjects")
            os.makedirs(root, exist_ok=True)
            for ck in VALID_COURSE_KEYS:
                os.makedirs(os.path.join(root, ck), exist_ok=True)
        except Exception as e:
            print(f"⚠️  Chatbot: Could not create subject Chroma folders: {e}")

    def _subject_chroma_path(self, course_key: str) -> str:
        return os.path.join(self.chroma_db_path, "subjects", course_key)

    def _course_material_chunk_count(self, course_key: str) -> int:
        path = self._subject_chroma_path(course_key)
        if not os.path.isdir(path) or chromadb is None:
            return 0
        try:
            if Settings:
                cli = chromadb.PersistentClient(
                    path=path,
                    settings=Settings(anonymized_telemetry=False, allow_reset=True),
                )
            else:
                cli = chromadb.PersistentClient(path=path)
            coll = cli.get_collection(COURSE_MATERIALS_COLLECTION)
            return coll.count()
        except Exception:
            return 0

    def _get_course_material_store(self, course_key: str):
        """LangChain Chroma store for one course; separate persist_directory from main timetable/schedules DB."""
        if course_key not in VALID_COURSE_KEYS:
            return None
        if not self.embeddings or Chroma is None or Document is None:
            return None
        if course_key in self._course_material_store_cache:
            return self._course_material_store_cache[course_key]
        path = self._subject_chroma_path(course_key)
        os.makedirs(path, exist_ok=True)
        try:
            try:
                store = Chroma(
                    collection_name=COURSE_MATERIALS_COLLECTION,
                    embedding_function=self.embeddings,
                    persist_directory=path,
                )
            except TypeError:
                store = Chroma(
                    collection_name=COURSE_MATERIALS_COLLECTION,
                    embedding=self.embeddings,
                    persist_directory=path,
                )
        except Exception as e:
            print(f"⚠️  Chatbot: Could not open course-materials store for {course_key}: {e}")
            return None
        self._course_material_store_cache[course_key] = store
        return store

    def _extract_course_material_text(self, file_path: str) -> str:
        ext = Path(file_path).suffix.lower()
        if ext == ".pdf":
            return self._extract_text_from_pdf_file(file_path)
        if ext == ".docx" and DocxDocument is not None:
            try:
                doc = DocxDocument(file_path)
                return "\n".join(p.text for p in doc.paragraphs if (p.text or "").strip()).strip()
            except Exception as e:
                print(f"⚠️  Chatbot: DOCX extract failed: {e}")
                return ""
        if ext == ".pptx" and PptxPresentation is not None:
            try:
                prs = PptxPresentation(file_path)
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and (shape.text or "").strip():
                            parts.append(shape.text.strip())
                return "\n".join(parts).strip()
            except Exception as e:
                print(f"⚠️  Chatbot: PPTX extract failed: {e}")
                return ""
        return ""

    def ingest_course_material_file(
        self,
        course_key: str,
        file_path: str,
        original_filename: str,
        storage_id: str,
    ) -> Dict:
        """Chunk, embed, and persist vectors for one uploaded file under the course subject folder."""
        out = {"success": False, "chunks": 0, "error": None}
        if course_key not in VALID_COURSE_KEYS:
            out["error"] = "invalid_course"
            return out
        text = self._extract_course_material_text(file_path)
        if not (text or "").strip():
            out["error"] = "no_text_extracted"
            return out
        if not RecursiveCharacterTextSplitter:
            out["error"] = "no_splitter"
            return out
        store = self._get_course_material_store(course_key)
        if not store:
            out["error"] = "no_vector_store"
            return out
        try:
            splitter = RecursiveCharacterTextSplitter(chunk_size=900, chunk_overlap=120)
            chunks = splitter.split_text(text)
            if not chunks:
                out["error"] = "empty_chunks"
                return out
            metadatas = [
                {
                    "source": original_filename,
                    "storage_id": storage_id,
                    "course": course_key,
                }
                for _ in chunks
            ]
            docs = [
                Document(page_content=c, metadata=metadatas[i])
                for i, c in enumerate(chunks)
            ]
            ids = [f"{storage_id}_{i}" for i in range(len(docs))]
            try:
                store.add_documents(docs, ids=ids)
            except TypeError:
                store.add_documents(docs)
            out["success"] = True
            out["chunks"] = len(docs)
            return out
        except Exception as e:
            print(f"⚠️  Chatbot: ingest_course_material_file failed: {e}")
            traceback.print_exc()
            out["error"] = str(e)
            return out

    def delete_course_material_vectors(self, course_key: str, storage_id: str) -> bool:
        """Remove all chunks for one upload from the course subject vector store."""
        if course_key not in VALID_COURSE_KEYS or not storage_id:
            return False
        store = self._get_course_material_store(course_key)
        if not store:
            return False
        try:
            coll = getattr(store, "_collection", None)
            if coll is None:
                return False
            coll.delete(where={"storage_id": storage_id})
            self._course_material_store_cache.pop(course_key, None)
            return True
        except Exception as e:
            print(f"⚠️  Chatbot: delete_course_material_vectors failed: {e}")
            return False

    def _load_course_material_manifest_entries(self, course_key: str) -> List[Dict]:
        path = os.path.join(self._course_materials_upload_root, course_key, "manifest.json")
        if not os.path.isfile(path):
            return []
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
            return data if isinstance(data, list) else []
        except Exception:
            return []

    def _aggregate_full_course_material_texts(self, course_key: str) -> Tuple[str, List[str]]:
        """Read all uploaded course files from disk (full text), for LLM when vector match is weak."""
        entries = self._load_course_material_manifest_entries(course_key)
        entries.sort(key=lambda e: (e.get("uploaded_at") or ""))
        parts: List[str] = []
        names: List[str] = []
        folder = os.path.join(self._course_materials_upload_root, course_key)
        for e in entries:
            save_as = e.get("saved_as") or ""
            orig = e.get("original_filename") or save_as
            fp = os.path.join(folder, save_as)
            if not save_as or not os.path.isfile(fp):
                continue
            txt = self._extract_course_material_text(fp)
            if not (txt or "").strip():
                continue
            parts.append(f"=== File: {orig} ===\n{txt.strip()}")
            if orig not in names:
                names.append(orig)
        return "\n\n".join(parts).strip(), names

    def answer_course_materials_question(
        self,
        question: str,
        user_course: Optional[str],
        db=None,
        user_email: Optional[str] = None,
        user_role: Optional[str] = None,
        student_mode: bool = False,
    ) -> Dict:
        """
        /course intent: match query to embeddings; if match is strong, use retrieved chunks for the LLM.
        If match is weak, skip chunk context and pass concatenated full file text (truncated) to the LLM.
        Always labels 'uploads' vs 'AI model' when any upload context exists; otherwise general LLM only.
        """
        if not question.strip():
            return {"answer": "Please ask a question.", "query_type": "course", "sources": []}
        if not user_course or user_course not in VALID_COURSE_KEYS:
            return {
                "answer": "Your profile has no course assigned. Please contact admin.",
                "query_type": "course",
                "sources": [],
            }
        if not self._groq_api_key:
            return {
                "answer": "Course Q&A needs GROQ_API_KEY to be configured.",
                "query_type": "course",
                "sources": [],
            }

        role = (user_role or "").strip().lower()
        uploaded_topics = self._load_user_attendance_topic_names(
            db=db,
            user_email=user_email or "",
            user_role=role,
            user_course=user_course,
        ) if db is not None and role in {"admin", "student"} else []
        manifest_entries = self._load_course_material_manifest_entries(user_course)
        uploaded_material_names = []
        for entry in manifest_entries:
            name = (entry.get("original_filename") or entry.get("saved_as") or "").strip()
            if name and name not in uploaded_material_names:
                uploaded_material_names.append(name)

        if self._is_topic_listing_question(question):
            topic_text = self._format_topic_list_response(uploaded_topics, user_role=role if role in {"admin", "student"} else "student")
            if uploaded_material_names:
                materials_block = "\n".join([f"- {n}" for n in uploaded_material_names[:20]])
                topic_text += f"\n\nUploaded course files you can ask from:\n{materials_block}"
            return {
                "answer": topic_text,
                "query_type": "course_topics_list",
                "sources": ["attendance_topics", "course_materials"],
            }

        store = self._get_course_material_store(user_course)
        pairs = []
        best_distance: Optional[float] = None
        chunk_context = ""
        source_files: List[str] = []

        if store:
            pairs = self._similarity_search_with_score_safe(store, question, k=8)
            for doc, sc in pairs:
                pc = (getattr(doc, "page_content", None) or "").strip()
                if not pc:
                    continue
                if sc is not None and isinstance(sc, (int, float)):
                    if best_distance is None or float(sc) < best_distance:
                        best_distance = float(sc)
                chunk_context += pc + "\n\n"
                src = (doc.metadata or {}).get("source")
                if src and src not in source_files:
                    source_files.append(src)

        chunk_context = (chunk_context or "").strip()
        use_chunk_rag = bool(
            chunk_context
            and best_distance is not None
            and best_distance <= COURSE_MATERIAL_RAG_MAX_DISTANCE
        )

        if use_chunk_rag:
            context = chunk_context[:12000]
            mode = "course_rag_chunks"
        else:
            full_text, file_names = self._aggregate_full_course_material_texts(user_course)
            context = full_text[:COURSE_MATERIAL_FULL_CONTEXT_MAX_CHARS]
            mode = "course_full_files_fallback"
            if file_names:
                source_files = file_names

        topics_only_scope = bool(uploaded_topics) and self._is_related_to_uploaded_topics(question, uploaded_topics)
        files_only_scope = False
        if best_distance is not None:
            files_only_scope = best_distance <= COURSE_MATERIAL_RELATED_MAX_DISTANCE
        elif context.strip():
            files_only_scope = self._has_min_context_overlap(question, context)

        if not context.strip() and not uploaded_topics:
            return {
                "answer": "No course files or attendance topics are uploaded yet for your course. Please ask your admin to upload them before using /course.",
                "query_type": "course_no_uploads",
                "sources": [],
            }

        if not files_only_scope and not topics_only_scope:
            return {
                "answer": "This question is outside your uploaded course files and attendance topics.",
                "query_type": "course_out_of_scope",
                "sources": ["course_materials", "attendance_topics"],
            }

        topics_context = ", ".join(uploaded_topics[:80]) if uploaded_topics else "None uploaded"
        materials_context = ", ".join(uploaded_material_names[:40]) if uploaded_material_names else "None uploaded"
        solo = self._query_general_llm(question)
        solo = self._sanitize_chatbot_answer(solo) or ""

        dual_prompt = (
            "You are a strict course assistant.\n"
            "Respond using EXACTLY these two markdown sections (with these headings):\n\n"
            "### From uploaded course materials\n"
            "Answer using ONLY the provided course uploads context and uploaded attendance topic names.\n"
            "If the upload context does not contain enough detail, clearly say what is missing.\n\n"
            "### General explanation (AI model)\n"
            "Give a short teaching explanation from general knowledge.\n\n"
            "Scope rule: the user question is already in-scope to uploaded course files/topics.\n"
            "Do not change the topic.\n\n"
            f"Uploaded attendance topic names: {topics_context}\n"
            f"Uploaded course file names: {materials_context}\n\n"
            f"--- Course uploads context ---\n{context}\n--- end context ---\n\n"
            f"Optional model-only draft: {solo}\n\n"
            f"User question: {question}\n"
        )
        scoped_answer = self._complete_prompt(dual_prompt, max_tokens=1200)
        scoped_answer = self._sanitize_chatbot_answer(scoped_answer) or solo or context[:500]
        return {
            "answer": scoped_answer,
            "query_type": mode,
            "sources": source_files or (["attendance_topics"] if uploaded_topics else ["course_materials"]),
            "retrieval_distance": best_distance,
        }

    def load_attendance_topics(self, db, user_email: str, user_role: str, user_course: Optional[str] = None):
        """Load attendance topics from MongoDB"""
        try:
            if not self.embeddings or Document is None:
                return False
            
            documents = []
            
            if user_role == 'admin' and user_course:
                # Admin (teacher): Get topics from their attendance_topics array (stored in `admin` collection)
                admin_collection = db.admin
                admin = admin_collection.find_one({"email": user_email})
                if admin and 'attendance_topics' in admin:
                    for entry in admin['attendance_topics']:
                        date = entry.get('date', '')
                        topics = entry.get('topics', {})
                        course = entry.get('course', '')
                        
                        # Format topics into text
                        topic_texts = []
                        for course_key, topic_list in topics.items():
                            if isinstance(topic_list, list):
                                topic_texts.extend(topic_list)
                            elif isinstance(topic_list, str):
                                topic_texts.append(topic_list)
                        
                        if topic_texts:
                            text = f"On {date}, topics covered in {course}: {', '.join(topic_texts)}"
                            documents.append(Document(
                                page_content=text,
                                metadata={"source": "attendance", "date": date, "course": course}
                            ))
            
            elif user_role == 'student':
                # Student: Get topics from their attendance records
                students_collection = None
                if user_course:
                    from database import get_course_collection
                    try:
                        students_collection = get_course_collection(db, user_course)
                    except:
                        pass
                
                if students_collection is not None:
                    student = students_collection.find_one({"email": user_email})
                    if student and 'attendance' in student:
                        for record in student['attendance']:
                            date = record.get('date', '')
                            topics = record.get('topics', {})
                            
                            topic_texts = []
                            for course_key, topic_list in topics.items():
                                if isinstance(topic_list, list):
                                    topic_texts.extend(topic_list)
                                elif isinstance(topic_list, str):
                                    topic_texts.append(topic_list)
                            
                            if topic_texts:
                                text = f"On {date}, topics covered: {', '.join(topic_texts)}"
                                documents.append(Document(
                                    page_content=text,
                                    metadata={"source": "attendance", "date": date}
                                ))
            
            if documents and RecursiveCharacterTextSplitter and Chroma:
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=500,
                    chunk_overlap=50
                )
                texts = text_splitter.split_documents(documents)
                
                # Use a user-scoped collection so one user's attendance data never leaks to another.
                collection_name = self._build_attendance_collection_name(
                    user_email=user_email,
                    user_role=user_role,
                    user_course=user_course,
                )
                try:
                    self.chroma_client.delete_collection(name=collection_name)
                except Exception:
                    pass

                self.chroma_client.create_collection(
                    name=collection_name,
                    metadata={"description": "Attendance topics embeddings for chatbot (scoped)"}
                )

                self.attendance_store = Chroma.from_documents(
                    documents=texts,
                    embedding=self.embeddings,
                    client=self.chroma_client,
                    collection_name=collection_name,
                    persist_directory=self.chroma_db_path
                )
                print(f"✅ Chatbot: Loaded {len(documents)} attendance topic documents into ChromaDB")
                return True
            else:
                if not documents:
                    print("⚠️  Chatbot: No attendance topics found")
                return False
        except Exception as e:
            print(f"⚠️  Chatbot: Error loading attendance topics: {e}")
            traceback.print_exc()
            return False
    
    def classify_query(self, query: str) -> str:
        """Classify query type: timetable, schedule, attendance, or out_of_scope"""
        query_lower = query.lower()
        
        # Timetable keywords
        timetable_keywords = [
            'class', 'where is my class', 'what class', 'when is class', 
            'timetable', 'schedule', 'time', 'period', 'subject', 'teacher',
            'room', 'location', 'now', 'current', 'today'
        ]
        
        # Schedule keywords
        schedule_keywords = [
            'schedule', 'what topics', 'what did we study', 'what did we read',
            'topics covered', 'syllabus', 'curriculum'
        ]
        
        # Attendance keywords
        attendance_keywords = [
            'attendance', 'topics', 'what topics', 'studied', 'covered',
            'learned', 'read', 'taught'
        ]
        
        timetable_score = sum(1 for kw in timetable_keywords if kw in query_lower)
        schedule_score = sum(1 for kw in schedule_keywords if kw in query_lower)
        attendance_score = sum(1 for kw in attendance_keywords if kw in query_lower)
        
        # Check for time/date related queries (likely timetable)
        if any(word in query_lower for word in ['now', 'current', 'today', 'when', 'what time']):
            timetable_score += 2
        
        # Determine query type
        max_score = max(timetable_score, schedule_score, attendance_score)
        
        if max_score == 0:
            return "out_of_scope"
        elif timetable_score == max_score:
            return "timetable"
        elif attendance_score == max_score:
            return "attendance"
        else:
            return "schedule"

    def _is_greeting_query(self, query: str) -> bool:
        """Detect simple greeting/small-talk queries (whole-message or short phrases)."""
        query_lower = (query or "").lower().strip()
        if not query_lower:
            return False
        # Short standalone greetings (no academic keywords)
        if re.match(r"^(hi|hello|hey|good morning|good afternoon|good evening)(\s+there)?[\s!.]*$", query_lower):
            return True
        greeting_patterns = [
            r"\bhow are you\b", r"\bhow's it going\b", r"\bwhat's up\b", r"\bhow do you do\b",
            r"\bgood morning\b", r"\bgood evening\b", r"\bgood afternoon\b",
        ]
        return any(re.search(pattern, query_lower) for pattern in greeting_patterns)

    def _respond_greeting(self, user_name: str) -> str:
        return f"Hello {user_name}! I am doing well, thank you for asking. How can I help you today?"

    def _is_admin_teach_topics_question(self, question: str) -> bool:
        """Detect admin questions like: 'How many topics I taught...'."""
        q = (question or "").lower()
        # Intents
        has_how_many = any(x in q for x in ["how many", "count", "number of"])
        has_teach = any(x in q for x in ["teach", "taught", "taught them"])
        has_topics = "topic" in q or "topics" in q
        has_schedule = any(x in q for x in ["schedule", "schedules", "scedu", "sched", "scedual"])
        # Keep this intent conservative so it doesn't override normal timetable/schedule Q&A
        return has_how_many and has_teach and has_topics and (has_schedule or "class" in q)

    def _is_admin_topics_query(self, question: str) -> bool:
        """
        Broader admin topics intent:
        supports phrases like 'what topics I taught to classes' even without 'how many'.
        """
        q = (question or "").lower()
        has_topics = any(x in q for x in ["topic", "topics", "covered", "syllabus"])
        has_teach_or_class = any(x in q for x in ["teach", "taught", "class", "classes", "students"])
        # avoid intercepting pure timetable location queries
        timetable_only = any(x in q for x in ["where is my class", "room", "location", "lecture"]) and not has_topics
        return has_topics and has_teach_or_class and not timetable_only

    def _extract_topics_from_context(self, context: str) -> List[str]:
        """
        Extract topic names from stored context.
        Attendance context is structured as: "topics covered in {course}: Topic1, Topic2, ..."
        """
        if not context:
            return []

        patterns = [
            # Attendance-style
            r"topics covered in [^:]+:\s*(.+?)(?:\r?\n|$)",
            r"topics covered:\s*(.+?)(?:\r?\n|$)",
            # Sometimes schedules context may contain "topics: A, B"
            r"topics\s*[:\-]\s*(.+?)(?:\r?\n|$)",
            r"topic\s*[:\-]\s*(.+?)(?:\r?\n|$)",
        ]

        found: List[str] = []
        for pat in patterns:
            for m in re.finditer(pat, context, flags=re.IGNORECASE | re.DOTALL):
                topics_str = (m.group(1) or "").strip()
                if not topics_str:
                    continue

                # Prefer comma split; otherwise split on common separators.
                if "," in topics_str:
                    parts = [p.strip() for p in topics_str.split(",")]
                else:
                    parts = [p.strip() for p in re.split(r"\band\b|;|\||/|\u2013|-", topics_str)]

                for p in parts:
                    cleaned = p.strip().strip('"').strip("'").strip(".").strip()
                    if not cleaned:
                        continue
                    # Avoid capturing "On 2026-..." or "Answer:" type strings
                    if re.match(r"^\d{4}-\d{2}-\d{2}\b", cleaned):
                        continue
                    found.append(cleaned)

        # De-duplicate while preserving order
        seen = set()
        unique = []
        for t in found:
            k = t.lower()
            if k in seen:
                continue
            seen.add(k)
            unique.append(t)
        return unique

    def _load_user_attendance_topic_names(self, db, user_email: str, user_role: str, user_course: Optional[str]) -> List[str]:
        """Fetch normalized topic names from attendance data for current user scope."""
        topics: List[str] = []
        try:
            if db is None:
                return topics

            if user_role == "admin" and user_course:
                admin = db.admin.find_one({"email": user_email})
                entries = (admin or {}).get("attendance_topics", [])
                for entry in entries:
                    topic_obj = entry.get("topics", {})
                    for _, v in (topic_obj or {}).items():
                        if isinstance(v, list):
                            topics.extend([str(x).strip() for x in v if str(x).strip()])
                        elif isinstance(v, str) and v.strip():
                            topics.append(v.strip())

            elif user_role == "student" and user_course:
                from database import get_course_collection
                students_collection = get_course_collection(db, user_course)
                student = students_collection.find_one({"email": user_email}) if students_collection is not None else None
                for record in (student or {}).get("attendance", []):
                    topic_obj = record.get("topics", {})
                    for _, v in (topic_obj or {}).items():
                        if isinstance(v, list):
                            topics.extend([str(x).strip() for x in v if str(x).strip()])
                        elif isinstance(v, str) and v.strip():
                            topics.append(v.strip())
        except Exception as e:
            print(f"⚠️  Chatbot: Failed to load attendance topic names: {e}")

        # Normalize and de-duplicate
        clean = []
        seen = set()
        for t in topics:
            n = re.sub(r"\s+", " ", t.lower()).strip()
            if not n or n in seen:
                continue
            seen.add(n)
            clean.append(t.strip())
        return clean

    def _is_related_to_uploaded_topics(self, question: str, uploaded_topics: List[str]) -> bool:
        """
        Heuristic relation check: if question references a topic token or close variation,
        treat it as in-scope for student follow-up.
        """
        q = (question or "").lower()
        q_norm = re.sub(r"[^a-z0-9\s]", " ", q)
        q_norm = re.sub(r"\s+", " ", q_norm).strip()

        # Canonical aliases for common CS terms
        if re.search(r"\bdsa\b", q_norm):
            q_norm += " data structure data structures algorithms"
        if "datastructure" in q_norm:
            q_norm += " data structure"
        if not q or not uploaded_topics:
            return False

        # direct phrase overlap
        for topic in uploaded_topics:
            t = topic.lower().strip()
            if not t:
                continue
            t_norm = re.sub(r"[^a-z0-9\s]", " ", t)
            t_norm = re.sub(r"\s+", " ", t_norm).strip()
            if t_norm in q_norm:
                return True
            # token overlap with singular/plural relax
            t_tokens = [tok for tok in re.split(r"[^a-z0-9]+", t_norm) if tok and len(tok) > 2]
            for tok in t_tokens:
                if tok in q_norm or (tok.endswith("s") and tok[:-1] in q_norm) or (f"{tok}s" in q_norm):
                    return True

        # Data-structure concept bridging for common classroom topics
        bridges = {
            "array": ["1d", "2d", "3d", "multidimensional", "index", "contiguous", "types of array"],
            "linked list": ["node", "next pointer", "singly", "doubly", "circular"],
            "stack": ["lifo", "push", "pop"],
            "queue": ["fifo", "enqueue", "dequeue"],
            "tree": ["binary tree", "bst", "traversal"],
            "graph": ["bfs", "dfs", "vertex", "edge"],
            "data structure": ["dsa", "algorithm", "algorithms", "time complexity", "space complexity", "linear", "non linear"],
            "data sturucture": ["dsa", "algorithm", "algorithms", "time complexity", "space complexity", "linear", "non linear"],
        }
        for base, hints in bridges.items():
            if any(base in (t or "").lower() for t in uploaded_topics):
                if any(h in q_norm for h in hints):
                    return True
        return False

    def _has_min_context_overlap(self, question: str, context: str, min_hits: int = 2) -> bool:
        """
        Lightweight lexical scope guard for /course questions.
        Used when vector distance is unavailable/weak.
        """
        q = (question or "").lower()
        c = (context or "").lower()
        if not q or not c:
            return False
        q_tokens = [t for t in re.split(r"[^a-z0-9]+", q) if len(t) >= 4]
        if not q_tokens:
            return False
        hits = 0
        seen = set()
        for tok in q_tokens:
            if tok in seen:
                continue
            seen.add(tok)
            if tok in c:
                hits += 1
                if hits >= min_hits:
                    return True
        return False

    def _answer_student_related_topic_fallback(self, question: str, uploaded_topics: List[str]) -> Optional[str]:
        """
        If question is related to uploaded topic names but retrieval misses,
        answer with LLM in a constrained manner (topic-scoped teaching response).
        """
        if not self._groq_api_key or not uploaded_topics:
            return None
        prompt = (
            "You are an educational assistant for students.\n"
            "Only answer if the question is related to one of these uploaded class topics.\n"
            "If related, provide a concise teaching explanation with simple examples.\n"
            "If not related, say: 'This question is outside the uploaded class topics.'\n\n"
            f"Uploaded topics: {', '.join(uploaded_topics[:40])}\n"
            f"Student question: {question}\n\n"
            "Answer:"
        )
        return self._complete_prompt(prompt, max_tokens=450)

    def _answer_student_course_mode(self, question: str, uploaded_topics: List[str]) -> Optional[str]:
        """
        Student /course mode:
        - Always bound the answer to uploaded topics (and closely related helpers).
        - If not related, refuse with the exact message used elsewhere.
        """
        if not uploaded_topics:
            return None
        if not self._is_related_to_uploaded_topics(question, uploaded_topics):
            return "This question is outside the uploaded class topics."
        if not self._groq_api_key:
            return None

        topic_list = ", ".join(uploaded_topics[:60])
        prompt = (
            "You are an educational assistant for students.\n"
            "STRICT SCOPE RULES:\n"
            "- You may ONLY answer questions that are about the uploaded class topics list below, "
            "or very closely related helper concepts needed to understand them.\n"
            "- If the question is not in-scope, reply exactly:\n"
            "  This question is outside the uploaded class topics.\n"
            "- Do not answer anything outside the scope.\n\n"
            f"Uploaded class topics: {topic_list}\n\n"
            f"Student question: {question}\n\n"
            "Answer in a clear, short teaching style with 1 small example if helpful.\n"
            "Answer:"
        )
        return self._complete_prompt(prompt, max_tokens=520)

    def _is_topic_listing_question(self, question: str) -> bool:
        """Detect requests asking for taught/uploaded topic names."""
        q = (question or "").lower()
        asks_topic = any(x in q for x in ["topic", "topics", "covered", "syllabus"])
        asks_list = any(x in q for x in ["what", "which", "list", "show", "tell me", "name"])
        return asks_topic and asks_list

    def _format_topic_list_response(self, topics: List[str], user_role: str) -> str:
        if not topics:
            if user_role == "student":
                return "I could not find uploaded taught topics yet. Please ask admin to upload attendance topics first."
            return "I could not find taught topics yet. Please upload schedules/attendance topics first."
        ordered = sorted({t.strip() for t in topics if t and str(t).strip()}, key=lambda x: x.lower())
        preview = ordered[:40]
        body = "\n".join([f"- {t}" for t in preview])
        more = ""
        if len(ordered) > len(preview):
            more = f"\n...and {len(ordered) - len(preview)} more topic(s)."
        return f"These are the taught/uploaded topics:\n{body}{more}"

    def query_admin_taught_topics(self, question: str, user_course: Optional[str] = None) -> Optional[str]:
        """Admin-only: count and list topics taught, using schedules vector DB first."""
        try:
            topics: List[str] = []

            # 1) Schedules vector DB (as requested)
            if self.schedules_store:
                retriever = self.schedules_store.as_retriever(search_kwargs={"k": 10})
                docs = self._retriever_get_docs(retriever, question)
                if docs:
                    if user_course:
                        course_l = user_course.lower()
                        docs = [
                            d for d in docs
                            if course_l in str(getattr(d, "metadata", {}) or {}).lower()
                            or course_l in str(getattr(d, "metadata", {}) or {}).get("file", "").lower()
                            or course_l in str(getattr(d, "metadata", {}) or {}).get("source", "").lower()
                        ]
                    context = "\n\n".join([d.page_content for d in docs if getattr(d, "page_content", "")])
                    topics = self._extract_topics_from_context(context)

            # 2) Attendance topics fallback (more structured)
            if not topics and self.attendance_store:
                retriever = self.attendance_store.as_retriever(search_kwargs={"k": 10})
                docs = self._retriever_get_docs(retriever, question)
                if docs:
                    if user_course:
                        course_l = user_course.lower()
                        docs = [d for d in docs if course_l in str(d.metadata.get("course", "")).lower()]
                    context = "\n\n".join([d.page_content for d in docs if getattr(d, "page_content", "")])
                    topics = self._extract_topics_from_context(context)

            if topics:
                return "You taught {} topic(s) to the class:\n{}".format(len(topics), ", ".join(topics))

            # 3) LLM fallback (only if we can actually call Groq)
            if self._groq_api_key and (self.schedules_store or self.attendance_store):
                context_parts = []
                if self.schedules_store:
                    retriever = self.schedules_store.as_retriever(search_kwargs={"k": 10})
                    docs = self._retriever_get_docs(retriever, question)
                    context_parts.append("\n\n".join([d.page_content for d in docs if getattr(d, "page_content", "")]))
                if self.attendance_store:
                    retriever = self.attendance_store.as_retriever(search_kwargs={"k": 10})
                    docs = self._retriever_get_docs(retriever, question)
                    context_parts.append("\n\n".join([d.page_content for d in docs if getattr(d, "page_content", "")]))

                context = "\n\n".join([c for c in context_parts if c.strip()])
                if context.strip():
                    prompt = (
                        "You are an assistant for an education system.\n"
                        "Admin request: count unique topics the admin taught and list them.\n\n"
                        f"Admin question: {question}\n\n"
                        f"Context:\n{context}\n\n"
                        "Respond exactly in this format:\n"
                        "You taught X topic(s) to the class:\n"
                        "Topic1, Topic2, Topic3"
                    )
                    return self._complete_prompt(prompt, max_tokens=350)

            return None
        except Exception as e:
            print(f"⚠️  Chatbot: Error querying admin taught topics: {e}")
            traceback.print_exc()
            return None

    def _retriever_get_docs(self, retriever, question: str) -> List:
        """LangChain retriever API differs across versions."""
        try:
            return retriever.get_relevant_documents(question)
        except Exception:
            pass
        try:
            out = retriever.invoke(question)
            if isinstance(out, list):
                return out
            if hasattr(out, "documents"):
                return out.documents or []
        except Exception as e:
            print(f"⚠️  Chatbot: retriever invoke failed: {e}")
        return []

    def _similarity_search_with_score_safe(self, store, question: str, k: int = 4):
        """Return list of (doc, score). Score meaning depends on store; Chroma uses distance (lower is better)."""
        if store is None:
            return []
        try:
            if hasattr(store, "similarity_search_with_score"):
                return store.similarity_search_with_score(question, k=k) or []
        except Exception as e:
            print(f"⚠️  Chatbot: similarity_search_with_score failed: {e}")
        # Fallback: no scores available
        try:
            docs = store.similarity_search(question, k=k) if hasattr(store, "similarity_search") else []
            return [(d, None) for d in (docs or [])]
        except Exception:
            return []

    def _best_student_vector_match(self, question: str):
        """
        Vector-first routing for students.
        Returns (route, best_score) where route is one of: 'attendance'|'schedule'|'timetable' or (None, None).
        For Chroma scores, lower distance means better match.
        """
        candidates = []
        for route, store in [
            ("attendance", self.attendance_store),
            ("schedule", self.schedules_store),
            ("timetable", self.timetable_store),
        ]:
            if not store:
                continue
            pairs = self._similarity_search_with_score_safe(store, question, k=8)
            if not pairs:
                continue
            # pick best scored doc; if score missing, just mark as potential match
            best = None
            for doc, score in pairs:
                if not getattr(doc, "page_content", "").strip():
                    continue
                if score is None:
                    best = ("doc", None)
                    break
                if best is None or (isinstance(score, (int, float)) and score < best[1]):
                    best = ("doc", float(score))
            if best is not None:
                candidates.append((route, best[1]))

        if not candidates:
            return None, None

        # Prefer lowest distance if we have numeric scores; otherwise first available route
        numeric = [(r, s) for (r, s) in candidates if isinstance(s, (int, float))]
        if numeric:
            best_route, best_score = sorted(numeric, key=lambda x: x[1])[0]
            if best_score <= STUDENT_RELATED_RAG_MAX_DISTANCE:
                return best_route, best_score
            return None, best_score

        # No scores available; treat as match
        return candidates[0][0], None

    def _run_retrieval_llm(self, store, question: str, prompt_template: str, k: int = 3, extra_inputs: Optional[Dict] = None) -> Optional[str]:
        """Retrieve context and call Groq (LangChain or HTTP fallback)."""
        if not store:
            return None
        if not self._groq_api_key:
            print("⚠️  Chatbot: no GROQ_API_KEY — cannot generate answer from context")
            return None

        try:
            retriever = store.as_retriever(search_kwargs={"k": k})
            docs = self._retriever_get_docs(retriever, question)
            if not docs:
                return None

            context = "\n\n".join([doc.page_content for doc in docs if getattr(doc, "page_content", "")])
            if not context.strip():
                return None

            payload = {"context": context, "question": question}
            if extra_inputs:
                payload.update(extra_inputs)

            formatted_prompt = prompt_template.format(**payload)
            return self._complete_prompt(formatted_prompt, max_tokens=512)
        except Exception as e:
            print(f"⚠️  Chatbot: retrieval+llm error: {e}")
            traceback.print_exc()
            return None

    def _run_llm_with_context(self, context: str, question: str, prompt_template: str, extra_inputs: Optional[Dict] = None) -> Optional[str]:
        """
        Run LLM on already-retrieved context (avoids second retrieval mismatch).
        """
        if not self._groq_api_key:
            return None
        if not (context or "").strip():
            return None
        try:
            payload = {"context": context, "question": question}
            if extra_inputs:
                payload.update(extra_inputs)
            formatted_prompt = prompt_template.format(**payload)
            return self._complete_prompt(formatted_prompt, max_tokens=512)
        except Exception as e:
            print(f"⚠️  Chatbot: llm-with-context error: {e}")
            traceback.print_exc()
            return None

    def _load_latest_timetable_context_from_files(self) -> str:
        """
        Load timetable context directly from the latest uploaded timetable file.
        No embeddings/vectors are used for timetable chatbot answering.
        """
        try:
            base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
            # Support both project layouts seen in this codebase:
            # - Backend/uploads/timetables
            # - Backend/timetable/uploads
            uploads_candidates = [
                os.path.join(base_dir, "uploads", "timetables"),
                os.path.join(base_dir, "timetable", "uploads"),
            ]
            files = []
            for folder in uploads_candidates:
                if not folder or not os.path.exists(folder):
                    continue
                for ext in ("*.png", "*.jpg", "*.jpeg", "*.pdf"):
                    files.extend(glob.glob(os.path.join(folder, ext)))
            if not files:
                return ""

            latest = max(files, key=os.path.getmtime)
            cache_key = f"upload:{latest}:{int(os.path.getmtime(latest))}"
            if self._timetable_context_cache and self._timetable_context_cache_source == cache_key:
                return self._timetable_context_cache

            suffix = Path(latest).suffix.lower()
            if suffix == ".pdf":
                context = self._extract_text_from_pdf_file(latest)
            else:
                context = self._extract_text_from_image_file(latest)

            self._timetable_context_cache_source = cache_key
            self._timetable_context_cache = (context or "").strip()
            return self._timetable_context_cache
        except Exception as e:
            print(f"⚠️  Chatbot: Could not load timetable context from files: {e}")
            return ""

    def _extract_text_from_pdf_file(self, pdf_path: str) -> str:
        try:
            if PdfReader is None:
                return ""
            reader = PdfReader(pdf_path)
            pages = []
            for page in reader.pages:
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    pages.append(page_text)
            return "\n".join(pages).strip()
        except Exception as e:
            print(f"⚠️  Chatbot: PDF text extraction failed: {e}")
            return ""

    def _extract_text_from_image_file(self, image_path: str) -> str:
        try:
            if pytesseract is None or Image is None:
                return ""
            # OCR enhancement pipeline for timetable grids:
            # grayscale -> autocontrast -> sharpen -> adaptive threshold.
            img = Image.open(image_path).convert("L")
            img = ImageOps.autocontrast(img)
            img = img.filter(ImageFilter.SHARPEN)

            # Try multiple OCR configs and keep the best output.
            configs = [
                "--oem 3 --psm 6",   # uniform block
                "--oem 3 --psm 4",   # single column
                "--oem 3 --psm 11",  # sparse text
            ]
            best = ""
            for cfg in configs:
                try:
                    text = pytesseract.image_to_string(img, config=cfg) or ""
                    if len(text.strip()) > len(best.strip()):
                        best = text
                except Exception:
                    continue
            return (best or "").strip()
        except Exception as e:
            print(f"⚠️  Chatbot: Image OCR extraction failed: {e}")
            return ""

    def clear_timetable_context_cache(self):
        """Discard cached timetable file context after admin delete/upload changes."""
        self._timetable_context_cache = ""
        self._timetable_context_cache_source = ""

    def _expand_student_query(self, question: str) -> str:
        """
        Expand student academic queries with related terms so semantically related
        questions still match uploaded topic context.
        """
        q = (question or "").strip()
        if not q:
            return q
        q_lower = q.lower()

        concept_map = {
            "array": ["arrays", "types of arrays", "1d array", "2d array", "multidimensional array", "dynamic array", "list"],
            "stack": ["stack", "lifo", "push", "pop", "top"],
            "queue": ["queue", "fifo", "enqueue", "dequeue", "front", "rear"],
            "linked list": ["linked list", "singly linked list", "doubly linked list", "node", "pointer"],
            "tree": ["tree", "binary tree", "bst", "traversal", "root", "leaf"],
            "graph": ["graph", "vertices", "edges", "bfs", "dfs"],
            "sorting": ["sorting", "bubble sort", "selection sort", "insertion sort", "merge sort", "quick sort"],
            "searching": ["searching", "linear search", "binary search"],
            "data structure": ["data structure", "algorithm", "time complexity", "space complexity"],
            "dsa": ["data structure", "data structures", "algorithms", "time complexity", "space complexity"],
            "data sturucture": ["data structure", "data structures", "algorithms", "time complexity", "space complexity"],
        }

        expansions = []
        for key, values in concept_map.items():
            if key in q_lower or any(v in q_lower for v in values):
                expansions.extend(values[:4])

        if not expansions:
            return q

        # Keep original question first; append concise expansion tail.
        expansion_tail = ", ".join(dict.fromkeys(expansions))
        return f"{q}. Related concepts: {expansion_tail}"

    def _parse_time_token_to_minutes(self, token: str) -> Optional[int]:
        """Parse time tokens like '09:30', '9:30 AM', '2 PM' into minutes from midnight."""
        if not token:
            return None
        t = token.strip().upper().replace(".", "")

        # 24h format HH:MM
        m = re.match(r"^(\d{1,2}):(\d{2})$", t)
        if m:
            h = int(m.group(1))
            minute = int(m.group(2))
            if 0 <= h <= 23 and 0 <= minute <= 59:
                return h * 60 + minute
            return None

        # 12h format H(:MM)? AM/PM
        m = re.match(r"^(\d{1,2})(?::(\d{2}))?\s*(AM|PM)$", t)
        if m:
            h = int(m.group(1))
            minute = int(m.group(2) or "0")
            ap = m.group(3)
            if not (1 <= h <= 12 and 0 <= minute <= 59):
                return None
            if h == 12:
                h = 0
            if ap == "PM":
                h += 12
            return h * 60 + minute
        return None

    def _extract_time_ranges_from_text(self, text: str) -> List[Tuple[int, int]]:
        """
        Extract class time ranges from timetable text.
        Supported examples:
        - 09:00-10:00
        - 9:00 AM - 10:30 AM
        - 2 PM to 3 PM
        """
        if not text:
            return []
        ranges: List[Tuple[int, int]] = []

        pattern = re.compile(
            r"(\d{1,2}(?::\d{2})?\s*(?:AM|PM)?|\d{1,2}:\d{2})\s*(?:-|to)\s*(\d{1,2}(?::\d{2})?\s*(?:AM|PM)?|\d{1,2}:\d{2})",
            flags=re.IGNORECASE
        )
        for start_raw, end_raw in pattern.findall(text):
            start_min = self._parse_time_token_to_minutes(start_raw)
            end_min = self._parse_time_token_to_minutes(end_raw)
            if start_min is None or end_min is None:
                continue
            if end_min <= start_min:
                continue
            ranges.append((start_min, end_min))
        return ranges

    def _extract_time_slot_lines(self, context: str) -> List[Tuple[int, int, str]]:
        """
        Extract per-line slot text with parsed ranges so we can detect "free" slots.
        Returns tuples: (start_min, end_min, original_line).
        """
        if not context:
            return []
        slots: List[Tuple[int, int, str]] = []
        line_pattern = re.compile(
            r"(\d{1,2}(?::\d{2})?\s*(?:AM|PM)?|\d{1,2}:\d{2})\s*(?:-|to)\s*(\d{1,2}(?::\d{2})?\s*(?:AM|PM)?|\d{1,2}:\d{2})",
            flags=re.IGNORECASE
        )
        for raw_line in context.splitlines():
            line = (raw_line or "").strip()
            if not line:
                continue
            m = line_pattern.search(line)
            if not m:
                continue
            start_min = self._parse_time_token_to_minutes(m.group(1))
            end_min = self._parse_time_token_to_minutes(m.group(2))
            if start_min is None or end_min is None or end_min <= start_min:
                continue
            slots.append((start_min, end_min, line))
        return slots

    def _is_free_slot_line(self, line: str) -> bool:
        l = (line or "").lower()
        free_words = [
            "free", "off", "no class", "break", "recess", "empty",
            "self study", "self-study", "lunch"
        ]
        return any(w in l for w in free_words)

    def _is_requested_day_in_line(self, line: str, requested_day: str) -> bool:
        d = (requested_day or "").lower()
        if d in {"", "today", "tomorrow"}:
            return True
        return d in (line or "").lower()

    def _detect_free_slot_for_requested_time(self, context: str, requested_day: str, target_minutes: Optional[int]) -> Optional[bool]:
        """
        Return:
        - True => matching slot exists and is marked free
        - False => matching slot exists but not marked free
        - None => no matching slot lines found
        """
        if target_minutes is None:
            return None
        slot_lines = self._extract_time_slot_lines(context)
        if not slot_lines:
            return None
        matched = []
        for start_min, end_min, line in slot_lines:
            if start_min <= target_minutes <= end_min and self._is_requested_day_in_line(line, requested_day):
                matched.append(line)
        if not matched:
            return None
        return any(self._is_free_slot_line(line) for line in matched)

    def _find_slot_line_for_day_time(self, context: str, requested_day: str, target_minutes: Optional[int]) -> Optional[str]:
        """Return the best matching timetable line for requested day/time."""
        if target_minutes is None:
            return None
        slot_lines = self._extract_time_slot_lines(context)
        if not slot_lines:
            return None
        candidates = []
        for start_min, end_min, line in slot_lines:
            if start_min <= target_minutes <= end_min:
                if self._is_requested_day_in_line(line, requested_day):
                    candidates.append((0, line))  # exact day hit
                else:
                    # keep as weak fallback only if day isn't explicit in line
                    if requested_day in {"today", "tomorrow"}:
                        candidates.append((1, line))
        if not candidates:
            return None
        candidates.sort(key=lambda x: x[0])
        return candidates[0][1]

    def _is_current_class_query(self, question: str) -> bool:
        q = (question or "").lower()
        return any(x in q for x in ["where is my class", "my class now", "current class", "class now", "right now", "at this time", "now"])

    def _has_class_at_current_time(self, context: str) -> Optional[bool]:
        """
        Return:
        - True if a time range includes current time
        - False if ranges exist but none include current time
        - None if no parseable ranges found
        """
        ranges = self._extract_time_ranges_from_text(context)
        if not ranges:
            return None
        now = datetime.now()
        now_min = now.hour * 60 + now.minute
        for start_min, end_min in ranges:
            if start_min <= now_min <= end_min:
                return True
        return False

    def _minutes_to_display_time(self, minutes: int) -> str:
        h = (minutes // 60) % 24
        m = minutes % 60
        suffix = "AM" if h < 12 else "PM"
        h12 = h % 12
        if h12 == 0:
            h12 = 12
        return f"{h12}:{m:02d} {suffix}"

    def _next_time_slot_hint(self, context: str) -> Optional[str]:
        """
        When there is no class right now, provide the nearest upcoming slot from
        extracted timetable ranges to keep response helpful/positive.
        """
        ranges = self._extract_time_ranges_from_text(context)
        if not ranges:
            return None
        now = datetime.now()
        now_min = now.hour * 60 + now.minute
        future = sorted([r for r in ranges if r[0] > now_min], key=lambda x: x[0])
        if future:
            start_min, end_min = future[0]
            return f"The next class slot starts around {self._minutes_to_display_time(start_min)}."
        # If all slots passed for today, surface earliest slot as a general hint.
        first_start, _ = sorted(ranges, key=lambda x: x[0])[0]
        return f"Today’s class slots seem finished. The earliest slot in timetable is around {self._minutes_to_display_time(first_start)}."

    def _extract_requested_datetime_context(self, question: str) -> Dict:
        """
        Extract requested day/time from user question.
        Returns keys: requested_day, requested_time_text, requested_time_minutes.
        """
        q = (question or "").lower()
        out = {
            "requested_day": "today",
            "requested_time_text": None,
            "requested_time_minutes": None,
        }
        if "tomorrow" in q:
            out["requested_day"] = "tomorrow"
        else:
            weekday_aliases = ["monday", "tuesday", "wednesday", "thursday", "friday", "saturday", "sunday"]
            for d in weekday_aliases:
                if d in q:
                    out["requested_day"] = d
                    break

        # parse first explicit time token from question
        m = re.search(r"(\d{1,2}(?::\d{2})?\s*(?:am|pm)|\d{1,2}:\d{2})", q, flags=re.IGNORECASE)
        if m:
            time_text = m.group(1).strip()
            out["requested_time_text"] = time_text
            out["requested_time_minutes"] = self._parse_time_token_to_minutes(time_text)
        return out

    def _has_class_at_specific_time(self, context: str, target_minutes: Optional[int]) -> Optional[bool]:
        if target_minutes is None:
            return None
        ranges = self._extract_time_ranges_from_text(context)
        if not ranges:
            return None
        for start_min, end_min in ranges:
            if start_min <= target_minutes <= end_min:
                return True
        return False

    def _context_mentions_requested_day(self, context: str, requested_day: str) -> bool:
        """
        Check whether retrieved timetable context contains requested day evidence.
        Prevents false-negative 'no class' for a specific weekday when day rows were not retrieved.
        """
        c = (context or "").lower()
        d = (requested_day or "").lower()
        if not c or not d:
            return False
        if d in {"today", "tomorrow"}:
            # Relative-day text often won't appear literally in uploaded timetable.
            return True
        return d in c

    def _sanitize_chatbot_answer(self, text: Optional[str]) -> Optional[str]:
        if not text:
            return text
        cleaned = re.sub(r"[ \t]+", " ", str(text))
        cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
        # guard against odd partial model artifacts
        if len(cleaned) < 8:
            return None
        return cleaned

    def _is_negative_timetable_fallback(self, answer: Optional[str]) -> bool:
        a = (answer or "").lower()
        negative_markers = [
            "don't have enough timetable information",
            "do not have enough timetable information",
            "i don't have enough information for that specific day/time",
            "i dont have enough information for that specific day/time",
            "i don't have any information for that specific",
            "i dont have any information for that specific",
            "specific schedule",
            "specific day/time",
            "not enough timetable information",
        ]
        return any(m in a for m in negative_markers)

    def _query_general_llm(self, question: str) -> Optional[str]:
        """General response path for admin users (any question)."""
        if not self._groq_api_key:
            return None
        prompt = (
            "You are a polite educational assistant. "
            "Answer clearly and briefly.\n\n"
            f"Question: {question}\n\n"
            "Answer:"
        )
        return self._complete_prompt(prompt, max_tokens=600)

    def _is_admin_student_data_question(self, question: str) -> bool:
        """Detect admin intent to fetch student data/details."""
        q = (question or "").lower()
        has_student = any(x in q for x in ["student", "students", "class list", "class strength"])
        has_detail = any(
            x in q
            for x in [
                "detail", "details", "information", "info", "record", "records",
                "profile", "phone", "address", "email", "father", "who is",
                "tell me about", "show", "list", "give me"
            ]
        )
        return has_student and has_detail

    def _is_sensitive_credential_request(self, question: str) -> bool:
        """Detect sensitive credential requests."""
        q = (question or "").lower()
        sensitive_keywords = [
            "password", "passcode", "otp", "token", "jwt", "secret", "hash", "hashed",
            "api key", "apikey", "credential", "credentials"
        ]
        return any(k in q for k in sensitive_keywords)

    def _is_file_request(self, question: str) -> bool:
        q = (question or "").lower()
        file_words = [
            "download",
            "file",
            "pdf",
            "document",
            "notes",
            "lecture notes",
            "timetable file",
            "schedule file",
            "syllabus",
            "give me link",
            "open file",
            "attach",
            "link to",
            "link for",
            "get the file",
            "where is the file",
            "which file",
            "uploaded file",
            "course file",
            "book file",
            "slide",
            "slides",
        ]
        if any(w in q for w in file_words):
            return True
        file_hint = any(
            x in q
            for x in ["file", "pdf", "document", "slides", "ppt", "notes", "upload", "link", "book"]
        )
        if file_hint and any(x in q for x in ["give me", "send me", "where can i get"]):
            return True
        return False

    def _build_course_material_download_link(self, course_key: str, storage_id: str) -> str:
        return f"{BACKEND_PUBLIC_BASE_URL}/api/course-materials/download/{course_key}/{storage_id}"

    def query_course_material_download_links(self, question: str, user_course: Optional[str]) -> Optional[str]:
        """
        For /course file requests: find relevant uploads via vector search, return markdown download links.
        Falls back to listing all course uploads if retrieval returns nothing.
        """
        if not user_course or user_course not in VALID_COURSE_KEYS:
            return None

        store = self._get_course_material_store(user_course)
        items: List[Dict] = []
        seen_ids = set()

        if store:
            pairs = self._similarity_search_with_score_safe(store, question, k=12)
            for doc, _sc in pairs:
                meta = getattr(doc, "metadata", {}) or {}
                sid = meta.get("storage_id")
                name = meta.get("source") or "document"
                if not sid or sid in seen_ids:
                    continue
                seen_ids.add(sid)
                items.append(
                    {
                        "filename": name,
                        "url": self._build_course_material_download_link(user_course, sid),
                    }
                )

        if not items:
            for e in self._load_course_material_manifest_entries(user_course):
                sid = e.get("storage_id")
                name = e.get("original_filename") or "document"
                if not sid or sid in seen_ids:
                    continue
                seen_ids.add(sid)
                items.append(
                    {
                        "filename": name,
                        "url": self._build_course_material_download_link(user_course, sid),
                    }
                )

        if not items:
            return None

        if len(items) == 1:
            f0 = items[0]
            return (
                f"The most relevant upload for your request is:\n"
                f"[Download {f0['filename']}]({f0['url']})"
            )

        lines = [
            "These course files may contain what you need — tap a link to download (opens in the app with your login):"
        ]
        for it in items[:6]:
            lines.append(f"- [{it['filename']}]({it['url']})")
        return "\n".join(lines)

    def _is_student_db_forbidden_request(self, question: str) -> bool:
        """Student must not access Mongo student records or credentials."""
        q = (question or "").lower()
        db_words = [
            "database", "record", "records", "mongo", "mongodb", "profile", "student info",
            "details of", "details about", "all students", "grade", "marks", "result"
        ]
        asks_sensitive = self._is_sensitive_credential_request(question)
        asks_db = any(w in q for w in db_words)
        asks_other_people = any(w in q for w in ["other student", "others", "everyone", "all students"])
        return asks_sensitive or asks_db or asks_other_people

    def _security_refusal_student(self) -> str:
        return (
            "I cannot provide student database records, passwords, or other private details. "
            "Please visit the Admin Office for personal or official student data requests."
        )

    def _embeddings_filename_to_original_filename(self, embeddings_filename: str) -> Optional[str]:
        """
        Convert saved embeddings filename back to uploaded file name.
        Example: timetable_20260101_120000_cs_png_embeddings.txt -> timetable_20260101_120000_cs.png
        """
        if not embeddings_filename or not embeddings_filename.endswith("_embeddings.txt"):
            return None
        base = embeddings_filename[: -len("_embeddings.txt")]
        if "_" not in base:
            return None
        stem, ext = base.rsplit("_", 1)
        ext = ext.lower()
        if ext in {"png", "jpg", "jpeg", "pdf"}:
            return f"{stem}.{ext}"
        return None

    def _build_download_link(self, source: str, filename: str) -> Optional[str]:
        if not filename:
            return None
        if source == "timetable":
            return f"{BACKEND_PUBLIC_BASE_URL}/api/timetable/download-timetable/{filename}"
        if source == "schedule":
            return f"{BACKEND_PUBLIC_BASE_URL}/api/schedules/download/{filename}"
        return None

    def _extract_file_candidates_from_docs(self, docs: List) -> List[Dict]:
        candidates = []
        seen = set()
        for d in docs or []:
            meta = getattr(d, "metadata", {}) or {}
            source = (meta.get("source") or "").lower()
            raw_file = meta.get("file") or ""
            original_file = self._embeddings_filename_to_original_filename(raw_file)
            if not original_file:
                continue
            link = self._build_download_link(source, original_file)
            if not link:
                continue
            key = (source, original_file)
            if key in seen:
                continue
            seen.add(key)
            candidates.append({
                "source": source,
                "filename": original_file,
                "url": link,
            })
        return candidates

    def query_file_download_links(self, question: str, user_role: str) -> Optional[str]:
        """
        Find file links using Chroma metadata (timetable/schedules).
        Works for both admin and student without touching student DB.
        """
        docs = []
        if self.timetable_store:
            docs.extend(self._retriever_get_docs(self.timetable_store.as_retriever(search_kwargs={"k": 6}), question))
        if self.schedules_store:
            docs.extend(self._retriever_get_docs(self.schedules_store.as_retriever(search_kwargs={"k": 6}), question))

        candidates = self._extract_file_candidates_from_docs(docs)
        if not candidates:
            return None

        if len(candidates) == 1:
            f = candidates[0]
            return f"[Download {f['filename']}]({f['url']})"

        # Multiple matches: ask user to specify.
        lines = ["I found multiple files. Please specify which one you need:"]
        for item in candidates[:6]:
            lines.append(f"- [{item['filename']}]({item['url']})")
        return "\n".join(lines)

    def _extract_student_identifier(self, question: str) -> Optional[str]:
        """Extract a likely student identifier from prompt (email/id/name fragment)."""
        q = (question or "").strip()
        email_match = re.search(r"([a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})", q)
        if email_match:
            return email_match.group(1).strip()
        reg_match = re.search(r"\bclass\d{2}\d{3}\b", q, flags=re.IGNORECASE)
        if reg_match:
            return reg_match.group(0).strip()
        # Capture "details of <name>" / "about <name>" patterns.
        name_match = re.search(r"(?:details?\s+of|about|for)\s+([a-zA-Z][a-zA-Z\s]{1,60})", q, flags=re.IGNORECASE)
        if name_match:
            candidate = name_match.group(1).strip().rstrip("?.!,")
            if candidate:
                return candidate
        return None

    def _safe_student_projection(self, student_doc: Dict) -> Dict:
        """Return fields requested by admin policy (includes stored password value)."""
        return {
            "name": student_doc.get("name") or student_doc.get("full_name") or "Unknown",
            "email": student_doc.get("email", ""),
            "password": student_doc.get("password", ""),
            "student_id": student_doc.get("student_id", ""),
            "class": student_doc.get("class", ""),
            "section": student_doc.get("section", ""),
            "phone": student_doc.get("phone", ""),
            "father_name": student_doc.get("father_name", ""),
            "address": student_doc.get("address", ""),
            "past_school": student_doc.get("past_school", ""),
            "is_active": student_doc.get("is_active", True),
            "total_present": student_doc.get("total_present", 0),
            "total_absent": student_doc.get("total_absent", 0),
        }

    def _format_student_data_response(self, student_rows: List[Dict]) -> str:
        if not student_rows:
            return "No matching students were found in your course data."
        if len(student_rows) == 1:
            s = student_rows[0]
            return (
                f"Student details:\n"
                f"- Name: {s.get('name')}\n"
                f"- Email: {s.get('email')}\n"
                f"- Password (stored): {s.get('password')}\n"
                f"- Student ID: {s.get('student_id')}\n"
                f"- Class/Section: {s.get('class')} {s.get('section')}\n"
                f"- Phone: {s.get('phone')}\n"
                f"- Father Name: {s.get('father_name')}\n"
                f"- Address: {s.get('address')}\n"
                f"- Previous School: {s.get('past_school')}\n"
                f"- Attendance: Present {s.get('total_present')} / Absent {s.get('total_absent')}\n"
                f"- Active: {s.get('is_active')}"
            )

        lines = ["Matching students in your course:"]
        for idx, s in enumerate(student_rows[:ADMIN_STUDENT_QUERY_MAX_DOCS], start=1):
            lines.append(
                f"{idx}. {s.get('name')} | ID: {s.get('student_id')} | Email: {s.get('email')} | Class: {s.get('class')} {s.get('section')}"
            )
        return "\n".join(lines)

    def _extract_requested_student_fields(self, question: str) -> List[str]:
        """
        Detect whether admin asked for specific entities only.
        Returns normalized field list; empty means full profile response.
        """
        q = (question or "").lower()
        field_map = [
            (["email", "mail"], "email"),
            (["password", "passcode", "credential"], "password"),
            (["phone", "contact", "mobile"], "phone"),
            (["father", "father name", "guardian"], "father_name"),
            (["address", "location"], "address"),
            (["school", "past school", "previous school"], "past_school"),
            (["student id", "registration", "roll no", "roll number", "id"], "student_id"),
            (["class", "section"], "class_section"),
            (["attendance", "present", "absent"], "attendance"),
            (["name", "full name"], "name"),
            (["active", "status"], "is_active"),
        ]

        requested = []
        for keywords, field in field_map:
            if any(k in q for k in keywords):
                requested.append(field)

        # De-duplicate while preserving order.
        out = []
        seen = set()
        for item in requested:
            if item in seen:
                continue
            seen.add(item)
            out.append(item)
        return out

    def _format_single_field_response(self, student_rows: List[Dict], requested_fields: List[str]) -> str:
        """Return only requested entities for one/multiple students."""
        if not student_rows:
            return "No matching students were found in your course data."
        if not requested_fields:
            return self._format_student_data_response(student_rows)

        label_map = {
            "name": "Name",
            "email": "Email",
            "password": "Password (stored)",
            "phone": "Phone",
            "father_name": "Father Name",
            "address": "Address",
            "past_school": "Previous School",
            "student_id": "Student ID",
            "class_section": "Class/Section",
            "attendance": "Attendance",
            "is_active": "Active",
        }

        def value_for(field: str, row: Dict):
            if field == "class_section":
                return f"{row.get('class', '')} {row.get('section', '')}".strip()
            if field == "attendance":
                return f"Present {row.get('total_present', 0)} / Absent {row.get('total_absent', 0)}"
            return row.get(field, "")

        if len(student_rows) == 1:
            s = student_rows[0]
            lines = [f"Requested details for {s.get('name', 'Student')}:"]
            for field in requested_fields:
                lines.append(f"- {label_map.get(field, field)}: {value_for(field, s)}")
            return "\n".join(lines)

        lines = ["Requested details for matching students:"]
        for idx, s in enumerate(student_rows[:ADMIN_STUDENT_QUERY_MAX_DOCS], start=1):
            parts = [f"{idx}. {s.get('name', 'Student')}"]
            for field in requested_fields:
                parts.append(f"{label_map.get(field, field)}: {value_for(field, s)}")
            lines.append(" | ".join(parts))
        return "\n".join(lines)

    def query_admin_student_data(self, db, question: str, user_course: Optional[str]) -> Optional[str]:
        """Admin-only student DB lookup, restricted to admin's own course collection."""
        if db is None or not user_course:
            return None

        try:
            from database import get_course_collection
            students_collection = get_course_collection(db, user_course)
        except Exception as e:
            print(f"⚠️  Chatbot: Failed to resolve course collection for admin student query: {e}")
            return None

        identifier = self._extract_student_identifier(question)
        query = {}
        if identifier:
            # Try exact hits first
            query = {
                "$or": [
                    {"email": identifier},
                    {"student_id": identifier},
                    {"name": {"$regex": re.escape(identifier), "$options": "i"}},
                    {"full_name": {"$regex": re.escape(identifier), "$options": "i"}},
                ]
            }

        cursor = students_collection.find(query).limit(ADMIN_STUDENT_QUERY_MAX_DOCS)
        docs = list(cursor)
        if not docs and identifier:
            # Fallback broad search by any token.
            query = {
                "$or": [
                    {"name": {"$regex": re.escape(identifier), "$options": "i"}},
                    {"full_name": {"$regex": re.escape(identifier), "$options": "i"}},
                    {"email": {"$regex": re.escape(identifier), "$options": "i"}},
                    {"student_id": {"$regex": re.escape(identifier), "$options": "i"}},
                ]
            }
            docs = list(students_collection.find(query).limit(ADMIN_STUDENT_QUERY_MAX_DOCS))

        safe_rows = [self._safe_student_projection(d) for d in docs]
        requested_fields = self._extract_requested_student_fields(question)
        return self._format_single_field_response(safe_rows, requested_fields)
    
    def get_current_time_context(self) -> str:
        """Get current time and date context"""
        now = datetime.now()
        current_time = now.strftime("%H:%M")
        current_date = now.strftime("%Y-%m-%d")
        day_name = now.strftime("%A")
        return f"Current date: {current_date} ({day_name}), Current time: {current_time}"
    
    def query_timetable(self, question: str) -> Optional[str]:
        """Query timetable with time context"""
        try:
            # Add time context to question
            time_context = self.get_current_time_context()
            asked_dt = self._extract_requested_datetime_context(question)
            asked_day = asked_dt.get("requested_day") or "today"
            asked_time_text = asked_dt.get("requested_time_text")
            asked_time_minutes = asked_dt.get("requested_time_minutes")
            asked_context = f"Requested day: {asked_day}"
            if asked_time_text:
                asked_context += f", Requested time: {asked_time_text}"
            enhanced_question = f"{question}. {time_context}. {asked_context}"

            latest_timetable_path = self._get_latest_timetable_upload_path()
            if latest_timetable_path and Path(latest_timetable_path).suffix.lower() in {".png", ".jpg", ".jpeg"}:
                prompt = (
                    "You are reading the uploaded timetable image directly.\n"
                    "Answer clearly using timetable content.\n"
                    "If slot is free/break, explicitly say slot is free.\n"
                    "If asked day/time not present, say that clearly.\n\n"
                    f"Current Context: {time_context}. {asked_context}\n"
                    f"Question: {question}"
                )
                direct_answer = self._groq_http_complete_with_image(prompt=prompt, image_path=latest_timetable_path)
                if direct_answer:
                    return self._sanitize_chatbot_answer(direct_answer)
                return "I could not read the timetable image from the model right now. Please try again."

            # No image found for direct vision path; PDF direct-vision is not supported here.
            return "No timetable image is available for direct model analysis. Please upload timetable as image (PNG/JPG)."

            
        except Exception as e:
            print(f"⚠️  Error querying timetable: {e}")
            traceback.print_exc()
            return None
    
    def query_schedules(self, question: str) -> Optional[str]:
        """Query schedules"""
        if not self.schedules_store:
            return None
        
        try:
            # Create retrieval chain from ChromaDB
            retriever = self.schedules_store.as_retriever(search_kwargs={"k": 3})
            docs = self._retriever_get_docs(retriever, question)
            context = "\n\n".join([doc.page_content for doc in docs if getattr(doc, "page_content", "")])
            
            prompt_template = """You are a helpful assistant that answers questions about schedules and curriculum.
Use the following context to answer the question.

Context: {context}

Question: {question}

Provide a short, clear answer. If the information is not in the context, say you don't have that information.

Answer:"""
            
            raw = self._run_llm_with_context(
                context=context,
                question=question,
                prompt_template=prompt_template,
            )
            cleaned = self._sanitize_chatbot_answer(raw)
            if self._is_negative_timetable_fallback(cleaned) and context.strip():
                return "I found schedule data, but I could not map that exact request confidently yet. Please ask with course/day/time or file name."
            return cleaned
        except Exception as e:
            print(f"⚠️  Error querying schedules: {e}")
            return None
    
    def query_attendance_topics(self, question: str) -> Optional[str]:
        """Query attendance topics"""
        if not self.attendance_store:
            return None
        
        try:
            # Create retrieval chain from ChromaDB
            retriever = self.attendance_store.as_retriever(search_kwargs={"k": 5})
            
            prompt_template = """You are a helpful assistant that answers questions about topics that were covered in class.
Use the following context to answer the question. Only answer about topics that are mentioned in the context.

Context: {context}

Question: {question}

Provide a short, clear answer. If the topic is not in the context (meaning it hasn't been studied yet), politely apologize and say that topic hasn't been covered yet.

Answer:"""
            
            return self._run_retrieval_llm(
                store=self.attendance_store,
                question=question,
                prompt_template=prompt_template,
                k=5
            )
        except Exception as e:
            print(f"⚠️  Error querying attendance topics: {e}")
            return None
    
    def answer_question(
        self,
        question: str,
        user_email: str,
        user_role: str,
        user_course: Optional[str] = None,
        db=None,
        user_name: str = "User",
        intent: Optional[str] = None,
    ) -> Dict:
        """Main method to answer questions using RAG"""
        try:
            question = (question or "").strip()
            intent = (intent or "").strip().lower() or None
            if intent not in {None, "timetable", "course", "student", "ai"}:
                intent = None
            if not question:
                return {
                    "answer": "Please type a question.",
                    "query_type": "empty",
                    "sources": [],
                }

            # Greetings first — before DB/Chroma work so simple "hi" never fails on load/embeddings
            if self._is_greeting_query(question):
                return {
                    "answer": self._respond_greeting(user_name or "User"),
                    "query_type": "greeting",
                    "sources": [],
                }

            # Reload attendance topics (they change frequently)
            if db is not None:
                self.load_attendance_topics(db, user_email, user_role, user_course)

            # Student: vector-first matching (semantic). If we find a match in vector DB, answer from that context.
            if user_role == "student":
                expanded_student_question = self._expand_student_query(question)
                if self._is_student_db_forbidden_request(question):
                    return {
                        "answer": self._security_refusal_student(),
                        "query_type": "security_refusal",
                        "sources": [],
                    }

                if self._is_file_request(question) and intent != "course":
                    links_answer = self.query_file_download_links(
                        question=expanded_student_question, user_role=user_role
                    )
                    if links_answer:
                        return {
                            "answer": links_answer,
                            "query_type": "file_download",
                            "sources": ["chromadb_metadata"],
                        }

                # Intent override: timetable mode persists until user changes it
                if intent == "timetable":
                    timetable_answer = self.query_timetable(question)
                    if timetable_answer:
                        return {
                            "answer": timetable_answer,
                            "query_type": "timetable",
                            "sources": ["timetable_image_llm"],
                        }
                    return {
                        "answer": "I couldn't find a timetable upload yet. Please ask admin to upload the timetable first.",
                        "query_type": "timetable",
                        "sources": [],
                    }

                # Student /course mode: file download ask → links; else embeddings + LLM.
                if intent == "course":
                    if self._is_file_request(question):
                        cm_links = self.query_course_material_download_links(
                            question=expanded_student_question, user_course=user_course
                        )
                        if cm_links:
                            return {
                                "answer": cm_links,
                                "query_type": "file_download",
                                "sources": ["course_materials"],
                            }
                    cm = self.answer_course_materials_question(
                        question=question,
                        user_course=user_course,
                        db=db,
                        user_email=user_email,
                        user_role=user_role,
                        student_mode=True,
                    )
                    return {
                        "answer": cm.get("answer", ""),
                        "query_type": cm.get("query_type", "course"),
                        "sources": cm.get("sources", []),
                    }

                # Students cannot use admin-only student DB mode
                if intent == "student":
                    return {
                        "answer": "Student database mode is only available for admins. Please use /course or /timetable.",
                        "query_type": "security_refusal",
                        "sources": [],
                    }

                if intent == "ai":
                    return {
                        "answer": "AI mode is only available for admins. Use /course for questions that use your course uploads and the AI.",
                        "query_type": "security_refusal",
                        "sources": [],
                    }

                # Dynamic per-question context switch for students:
                # - timetable questions -> same timetable flow as admin
                # - topic-listing questions -> list uploaded taught topics
                student_query_type = self.classify_query(expanded_student_question)
                if student_query_type == "timetable":
                    timetable_answer = self.query_timetable(question)
                    if timetable_answer:
                        return {
                            "answer": timetable_answer,
                            "query_type": "timetable",
                            "sources": ["timetable_image_llm"],
                        }

                if student_query_type in {"attendance", "schedule"} and self._is_topic_listing_question(expanded_student_question):
                    uploaded_topics = self._load_user_attendance_topic_names(
                        db=db,
                        user_email=user_email,
                        user_role=user_role,
                        user_course=user_course,
                    )
                    return {
                        "answer": self._format_topic_list_response(uploaded_topics, user_role="student"),
                        "query_type": "student_topics_list",
                        "sources": ["attendance_topics"],
                    }

                route, best_score = self._best_student_vector_match(expanded_student_question)
                if route == "attendance":
                    ans = self.query_attendance_topics(expanded_student_question)
                    if ans:
                        return {"answer": ans, "query_type": "attendance", "sources": []}
                elif route == "schedule":
                    ans = self.query_schedules(expanded_student_question)
                    if ans:
                        return {"answer": ans, "query_type": "schedule", "sources": []}
                elif route == "timetable":
                    ans = self.query_timetable(expanded_student_question)
                    if ans:
                        return {"answer": ans, "query_type": "timetable", "sources": []}

                # Topic-related fallback:
                # If no direct vector match but question relates to uploaded attendance topics
                # (e.g., uploaded "array" and asked "what is 2D array"), still answer.
                uploaded_topics = self._load_user_attendance_topic_names(
                    db=db,
                    user_email=user_email,
                    user_role=user_role,
                    user_course=user_course,
                )
                if self._is_related_to_uploaded_topics(expanded_student_question, uploaded_topics):
                    related_answer = self._answer_student_related_topic_fallback(
                        question=question,
                        uploaded_topics=uploaded_topics,
                    )
                    if related_answer and "outside the uploaded class topics" not in related_answer.lower():
                        return {
                            "answer": related_answer,
                            "query_type": "related_topic_fallback",
                            "sources": ["attendance_topics"],
                        }

                # No semantic match -> student apology (admin didn't upload / topic not covered)
                return {
                    "answer": "I apologize, but I could not find this in the uploaded content from admin. Please ask about topics that have been uploaded or covered in class.",
                    "query_type": "no_match",
                    "sources": [],
                }

            # Admin intent override: explicit student DB mode
            if user_role == "admin" and intent == "student":
                admin_student_answer = self.query_admin_student_data(
                    db=db,
                    question=question,
                    user_course=user_course,
                )
                if admin_student_answer:
                    return {
                        "answer": admin_student_answer,
                        "query_type": "admin_student_data",
                        "sources": ["students_db"],
                    }
                return {
                    "answer": "I couldn't find any matching students for that query in your course data.",
                    "query_type": "admin_student_data",
                    "sources": ["students_db"],
                }

            # Admin intent override: free-form AI mode (ChatGPT-like)
            if user_role == "admin" and intent == "ai":
                general_answer = self._query_general_llm(question)
                return {
                    "answer": general_answer or "I apologize, but I could not generate an answer right now.",
                    "query_type": "ai",
                    "sources": [],
                }

            # Admin intent override: timetable mode
            if user_role == "admin" and intent == "timetable":
                timetable_answer = self.query_timetable(question)
                if timetable_answer:
                    return {
                        "answer": timetable_answer,
                        "query_type": "timetable",
                        "sources": ["timetable_image_llm"],
                    }
                return {
                    "answer": "I couldn't find a timetable upload yet. Please upload the timetable first.",
                    "query_type": "timetable",
                    "sources": [],
                }

            # Admin intent override: course mode (uploaded materials: file links, or embeddings + LLM)
            if user_role == "admin" and intent == "course":
                if self._is_file_request(question):
                    cm_links = self.query_course_material_download_links(question=question, user_course=user_course)
                    if cm_links:
                        return {
                            "answer": cm_links,
                            "query_type": "file_download",
                            "sources": ["course_materials"],
                        }
                cm = self.answer_course_materials_question(
                    question=question,
                    user_course=user_course,
                    db=db,
                    user_email=user_email,
                    user_role=user_role,
                    student_mode=False,
                )
                return {
                    "answer": cm.get("answer", ""),
                    "query_type": cm.get("query_type", "course"),
                    "sources": cm.get("sources", []),
                }

            # Admin-only: taught-topics intent (count + list)
            if user_role == "admin" and self._is_admin_teach_topics_question(question):
                admin_answer = self.query_admin_taught_topics(question=question, user_course=user_course)
                if admin_answer:
                    return {
                        "answer": admin_answer,
                        "query_type": "admin_teach_topics",
                        "sources": [],
                    }

            # Admin topic queries should dynamically switch to schedules/attendance context.
            if user_role == "admin" and self._is_admin_topics_query(question):
                admin_topics_answer = self.query_admin_taught_topics(question=question, user_course=user_course)
                if admin_topics_answer:
                    return {
                        "answer": admin_topics_answer,
                        "query_type": "admin_topics",
                        "sources": ["schedules", "attendance"],
                    }

            if user_role == "admin" and self._is_file_request(question) and intent != "course":
                links_answer = self.query_file_download_links(question=question, user_role=user_role)
                if links_answer:
                    return {
                        "answer": links_answer,
                        "query_type": "file_download",
                        "sources": ["chromadb_metadata"],
                    }

            # Admin-only: direct student data requests from DB (course-scoped, safe fields only).
            if user_role == "admin" and self._is_admin_student_data_question(question):
                admin_student_answer = self.query_admin_student_data(
                    db=db,
                    question=question,
                    user_course=user_course,
                )
                if admin_student_answer:
                    return {
                        "answer": admin_student_answer,
                        "query_type": "admin_student_data",
                        "sources": ["students_db"],
                    }

            # Classify query
            query_type = self.classify_query(question)

            # (student path handled above; admin continues below)

            # Admin can ask anything
            if user_role == "admin" and query_type == "out_of_scope":
                general_answer = self._query_general_llm(question)
                return {
                    "answer": general_answer or "I apologize, but I could not generate an answer right now.",
                    "query_type": "general",
                    "sources": []
                }
            
            # Route to appropriate handler
            answer = None
            sources = []
            
            if query_type == "timetable":
                answer = self.query_timetable(question)
            elif query_type == "schedule":
                answer = self.query_schedules(question)
            elif query_type == "attendance":
                answer = self.query_attendance_topics(question)
            
            if not answer and user_role == "student":
                return {
                    "answer": "I apologize, but I could not find this topic in the uploaded content from admin. Please ask about topics that have been uploaded or covered in class.",
                    "query_type": query_type,
                    "sources": []
                }

            if not answer and user_role == "admin":
                fallback_answer = self._query_general_llm(question)
                if fallback_answer:
                    answer = fallback_answer
                else:
                    return {
                        "answer": "I apologize, but I couldn't find relevant information to answer your question right now.",
                        "query_type": query_type,
                        "sources": []
                    }
            
            return {
                "answer": answer,
                "query_type": query_type,
                "sources": sources
            }
        except Exception as e:
            print(f"⚠️  Error answering question: {e}")
            traceback.print_exc()
            return {
                "answer": "I apologize, but I encountered an error processing your question. Please try again.",
                "query_type": "error",
                "sources": []
            }
